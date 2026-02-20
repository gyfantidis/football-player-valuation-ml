"""
DAMA Hackathon 2026 — 5-minute presentation builder.
Produces: report/presentation.pptx
Run: python3 report/build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # mid navy
GOLD      = RGBColor(0xE9, 0x4F, 0x37)   # red-orange highlight
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC8, 0xD8, 0xE8)   # light blue-grey
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
YELLOW    = RGBColor(0xF3, 0x9C, 0x12)
SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)

# ── Helper utilities ──────────────────────────────────────────────────────────

def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_box(slide, title, bullets, left, top, width, height,
                   title_size=20, bullet_size=16, bg=ACCENT):
    add_rect(slide, left, top, width, height, bg)
    add_text(slide, title, left + Inches(0.15), top + Inches(0.1),
             width - Inches(0.3), Inches(0.45),
             size=title_size, bold=True, color=GOLD)
    body = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.55),
        width - Inches(0.3), height - Inches(0.65)
    )
    body.word_wrap = True
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = b
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = LIGHT

def accent_bar(slide):
    """Horizontal gold bar at bottom."""
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), GOLD)

def slide_number(slide, n, total=10):
    add_text(slide, f"{n} / {total}",
             SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
             Inches(1.0), Inches(0.35),
             size=12, color=LIGHT, align=PP_ALIGN.RIGHT)

# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]   # completely blank


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)

# Left decorative stripe
add_rect(sl, 0, 0, Inches(0.25), SLIDE_H, GOLD)

# Title
add_text(sl,
    "Predictive Modelling &\nArchetype Discovery for\nFootball Player Valuation",
    Inches(0.55), Inches(1.0), Inches(8.5), Inches(3.2),
    size=40, bold=True, color=WHITE)

# Subtitle
add_text(sl,
    "A Comparative ML Study  |  DAMA Hackathon 2026",
    Inches(0.55), Inches(4.3), Inches(8.5), Inches(0.6),
    size=20, color=GOLD, bold=True)

# Meta
add_text(sl,
    "Dataset: FIFA Player Performance & Market Value  (n = 2,800)\n"
    "Tasks: Regression  ·  Classification  ·  Clustering  ·  Interpretability",
    Inches(0.55), Inches(5.1), Inches(9.0), Inches(1.0),
    size=16, color=LIGHT)

# Right visual accent — stat boxes
for i, (label, val) in enumerate([
    ("2,800", "Players"),
    ("5", "ML Models"),
    ("R² 0.958", "Best R²"),
]):
    bx_left = Inches(10.0)
    bx_top  = Inches(1.5 + i * 1.7)
    add_rect(sl, bx_left, bx_top, Inches(2.8), Inches(1.4), ACCENT)
    add_text(sl, val,  bx_left + Inches(0.15), bx_top + Inches(0.1),
             Inches(2.5), Inches(0.45), size=14, color=GOLD, bold=True)
    add_text(sl, label, bx_left + Inches(0.15), bx_top + Inches(0.55),
             Inches(2.5), Inches(0.7), size=26, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Dataset
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 2)

add_text(sl, "Problem & Dataset", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         size=32, bold=True, color=GOLD)
add_rect(sl, Inches(0.5), Inches(1.0), Inches(7.0), Inches(0.04), GOLD)

# Left: problem statement
add_bullet_box(sl, "Research Questions",
    ["RQ1  Can supervised models predict player market value?",
     "RQ2  Can classifiers identify transfer risk level?",
     "RQ3  Do performance features reveal player archetypes?"],
    Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.2), bullet_size=15)

# Right: dataset stats
add_bullet_box(sl, "Dataset Features",
    ["age, overall_rating, potential_rating",
     "goals, assists, minutes_played, matches_played",
     "contract_years_left, injury_prone",
     "position (9), nationality (8), club (7)"],
    Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.2), bullet_size=15)

# Target variable boxes
for i, (label, desc, col) in enumerate([
    ("market_value_million_eur", "Regression target  |  €0.7M – €180M  |  mean €90.6M", ACCENT),
    ("transfer_risk_level", "Classification target  |  Low 44.6%  ·  Medium 35.4%  ·  High 20.0%", ACCENT),
    ("FIFA Performance Value Index", "Constructed target (primary regression)  |  €0.5M – €181M", RGBColor(0x1A, 0x3A, 0x2E)),
]):
    add_rect(sl, Inches(0.5), Inches(3.6 + i * 1.1), Inches(12.3), Inches(0.9), col)
    add_text(sl, label, Inches(0.7), Inches(3.65 + i * 1.1),
             Inches(4.5), Inches(0.5), size=14, bold=True, color=GOLD)
    add_text(sl, desc, Inches(5.3), Inches(3.65 + i * 1.1),
             Inches(7.3), Inches(0.5), size=14, color=LIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — EDA: Critical Finding
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 3)

add_text(sl, "EDA: A Critical Data Finding", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         size=32, bold=True, color=GOLD)

# Big warning box
add_rect(sl, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.5), RGBColor(0x3A, 0x10, 0x10))
add_text(sl, "⚠  market_value_million_eur is uncorrelated with ALL features  (|r| < 0.03 for every variable)",
         Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.7),
         size=18, bold=True, color=YELLOW)
add_text(sl, "With n=2,800, the 95% CI on a zero correlation is ±0.037 — every feature lies within this band.",
         Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.55),
         size=14, color=LIGHT)

# Correlation table
headers = ["Feature", "|r| with market value", "Interpretation"]
rows = [
    ("overall_rating",   "0.013", "No signal"),
    ("age",              "0.014", "No signal"),
    ("goals",            "0.022", "No signal"),
    ("potential_rating", "0.006", "No signal"),
    ("contract_years_left", "0.012", "No signal"),
]
col_xs  = [Inches(0.5), Inches(4.5), Inches(9.0)]
col_ws  = [Inches(3.8), Inches(4.3), Inches(3.8)]
row_h   = Inches(0.5)

for ci, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
    add_rect(sl, cx, Inches(2.85), cw - Inches(0.05), row_h, RGBColor(0x10, 0x30, 0x50))
    add_text(sl, hdr, cx + Inches(0.1), Inches(2.9), cw, row_h - Inches(0.05),
             size=14, bold=True, color=GOLD)

for ri, (feat, corr, interp) in enumerate(rows):
    top = Inches(3.35) + ri * row_h
    bg  = ACCENT if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)
    for val, cx, cw in zip([feat, corr, interp], col_xs, col_ws):
        add_rect(sl, cx, top, cw - Inches(0.05), row_h - Inches(0.05), bg)
        add_text(sl, val, cx + Inches(0.1), top + Inches(0.05), cw, row_h,
                 size=13, color=WHITE if val != interp else RGBColor(0xFF, 0x88, 0x88))

# Conclusion box
add_rect(sl, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9), RGBColor(0x10, 0x30, 0x15))
add_text(sl,
    "→  Synthetic dataset with randomly assigned market values.  "
    "We document this finding and define a domain-informed FIFA Performance Value Index (FPVI) as primary target.",
    Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.7), size=14, color=GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methodology Overview
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 4)

add_text(sl, "Methodology", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         size=32, bold=True, color=GOLD)

steps = [
    ("1  Data", ["2,800 players · 13 raw features", "9 positions · 8 nations · 7 clubs"]),
    ("2  Engineering", ["goals/assists per 90 (clipped)", "rating_gap · age_rating_ratio", "expiring_soon · position_group", "FPVI target construction"]),
    ("3  Preprocessing", ["Ordinal + one-hot encoding", "StandardScaler (linear models)", "70 / 15 / 15 stratified split"]),
    ("4  Models", ["Ridge  ·  Random Forest", "Gradient Boosting", "XGBoost  ·  LightGBM", "Logistic Regression (cls)"]),
    ("5  Explain", ["SHAP TreeExplainer", "Global beeswarm + bar", "Local waterfall plots", "Dependence plots"]),
    ("6  Cluster", ["K-Means + silhouette k", "PCA 2D projection", "Radar profile charts"]),
]

box_w = Inches(2.05)
for i, (title, bullets) in enumerate(steps):
    bx = Inches(0.4) + i * (box_w + Inches(0.05))
    add_bullet_box(sl, title, bullets, bx, Inches(1.1), box_w, Inches(5.8),
                   title_size=16, bullet_size=12)

# Arrow between boxes (simple text arrows)
for i in range(len(steps) - 1):
    ax = Inches(0.4) + i * (box_w + Inches(0.05)) + box_w + Inches(0.0)
    add_text(sl, "→", ax - Inches(0.02), Inches(3.7), Inches(0.12), Inches(0.5),
             size=18, bold=True, color=GOLD)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FPVI Definition
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 5)

add_text(sl, "FIFA Performance Value Index (FPVI)", Inches(0.5), Inches(0.3),
         Inches(12), Inches(0.7), size=32, bold=True, color=GOLD)

# Formula box
add_rect(sl, Inches(0.5), Inches(1.1), Inches(12.3), Inches(2.4), RGBColor(0x0D, 0x1B, 0x2A))
add_text(sl,
    "age_factor  =  exp( −0.08 × max(0, age − 26)² )   clipped [0.1, 1.0]\n\n"
    "FPVI  =  rating_norm × 100 × age_factor\n"
    "       + pot_gap × 1.2 × age_factor\n"
    "       + goals_per_90 × 8  +  assists_per_90 × 5  +  ε",
    Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.2),
    size=18, color=WHITE, italic=False)

# Three principle boxes
principles = [
    ("Quality × Age Prime",
     "Rating decays exponentially after age 26.\n"
     "Peak contribution at age 24–27.", GREEN),
    ("Development Upside",
     "Young players with high potential gap\n"
     "(potential − overall) earn a premium.", YELLOW),
    ("On-pitch Production",
     "Goals and assists per 90 min add\n"
     "incremental value above base rating.", RGBColor(0x5D, 0xAD, 0xFF)),
]
for i, (title, body, col) in enumerate(principles):
    bx = Inches(0.5) + i * Inches(4.15)
    add_rect(sl, bx, Inches(3.7), Inches(3.95), Inches(2.5), ACCENT)
    add_rect(sl, bx, Inches(3.7), Inches(3.95), Inches(0.12), col)
    add_text(sl, title, bx + Inches(0.25), Inches(3.8),
             Inches(3.6), Inches(0.55), size=16, bold=True, color=col)
    add_text(sl, body, bx + Inches(0.25), Inches(4.45),
             Inches(3.6), Inches(1.5), size=14, color=LIGHT)

add_text(sl,
    "±10% Gaussian noise injected to prevent trivial learning   |   Range: €0.5M – €181M   |   Mean: €48.9M   |   Std: €34.2M",
    Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
    size=13, color=LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Regression Results
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 6)

add_text(sl, "Regression Results", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         size=32, bold=True, color=GOLD)

# Task A table (left)
add_text(sl, "Task A — Original Market Value (Negative Control)",
         Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=YELLOW)

mv_rows = [
    ("Ridge",             "54.6", "45.4", "−0.200"),
    ("Random Forest",     "55.2", "45.6", "−0.226"),
    ("Gradient Boosting", "56.9", "46.5", "−0.300"),
    ("XGBoost",           "57.6", "47.1", "−0.335"),
    ("LightGBM",          "57.9", "47.4", "−0.349"),
]
hdrs_a = ["Model", "RMSE", "MAE", "R²"]
col_xs_a = [Inches(0.5), Inches(3.0), Inches(4.2), Inches(5.2)]
col_ws_a = [Inches(2.4), Inches(1.1), Inches(0.9), Inches(1.0)]

for ci, (h, cx, cw) in enumerate(zip(hdrs_a, col_xs_a, col_ws_a)):
    add_rect(sl, cx, Inches(1.5), cw - Inches(0.05), Inches(0.45),
             RGBColor(0x10, 0x30, 0x50))
    add_text(sl, h, cx + Inches(0.08), Inches(1.55), cw, Inches(0.35),
             size=13, bold=True, color=GOLD)

for ri, row in enumerate(mv_rows):
    top = Inches(1.95) + ri * Inches(0.48)
    bg = ACCENT if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)
    for val, cx, cw in zip(row, col_xs_a, col_ws_a):
        add_rect(sl, cx, top, cw - Inches(0.05), Inches(0.43), bg)
        color = RGBColor(0xFF, 0x60, 0x60) if "−" in str(val) else WHITE
        add_text(sl, val, cx + Inches(0.08), top + Inches(0.05),
                 cw, Inches(0.35), size=12, color=color)

add_text(sl, "→ All R² < 0  confirms no signal in target",
         Inches(0.5), Inches(4.45), Inches(6.0), Inches(0.5),
         size=13, color=RGBColor(0xFF, 0x88, 0x88), italic=True)

# Task B table (right)
add_text(sl, "Task B — FIFA Performance Index  (Primary)",
         Inches(6.9), Inches(1.05), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=GREEN)

fp_rows = [
    ("LightGBM",          "6.87", "4.49", "0.958", True),
    ("XGBoost",           "6.96", "4.48", "0.957", False),
    ("Gradient Boosting", "7.04", "4.55", "0.956", False),
    ("Random Forest",     "9.81", "6.30", "0.915", False),
    ("Ridge",             "20.4", "13.8", "0.631", False),
]
col_xs_b = [Inches(6.9), Inches(9.4), Inches(10.55), Inches(11.6)]
col_ws_b = [Inches(2.4),  Inches(1.0),  Inches(0.95),  Inches(1.1)]
hdrs_b   = ["Model", "RMSE", "MAE", "R²"]

for ci, (h, cx, cw) in enumerate(zip(hdrs_b, col_xs_b, col_ws_b)):
    add_rect(sl, cx, Inches(1.5), cw - Inches(0.05), Inches(0.45),
             RGBColor(0x10, 0x30, 0x50))
    add_text(sl, h, cx + Inches(0.08), Inches(1.55), cw, Inches(0.35),
             size=13, bold=True, color=GOLD)

for ri, (name, rmse, mae, r2, best) in enumerate(fp_rows):
    top = Inches(1.95) + ri * Inches(0.48)
    bg = RGBColor(0x0D, 0x2A, 0x15) if best else (ACCENT if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E))
    for val, cx, cw in zip([name, rmse, mae, r2], col_xs_b, col_ws_b):
        add_rect(sl, cx, top, cw - Inches(0.05), Inches(0.43), bg)
        color = GREEN if best else WHITE
        add_text(sl, val, cx + Inches(0.08), top + Inches(0.05),
                 cw, Inches(0.35), size=12, bold=best, color=color)

add_text(sl, "★  LightGBM  R² = 0.958  |  5-fold CV R² = 0.968 ± 0.007",
         Inches(6.9), Inches(4.45), Inches(6.0), Inches(0.5),
         size=13, color=GREEN, bold=True)

# Key insight
add_rect(sl, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.0), RGBColor(0x10, 0x1A, 0x30))
add_text(sl,
    "Non-linearity matters:  Ridge R² = 0.63  vs  LightGBM R² = 0.96 — "
    "the age × rating interaction cannot be captured by a linear model without explicit feature crosses.",
    Inches(0.7), Inches(5.18), Inches(11.9), Inches(0.8), size=14, color=LIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SHAP Interpretability
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 7)

add_text(sl, "Model Interpretability — SHAP", Inches(0.5), Inches(0.3),
         Inches(12), Inches(0.7), size=32, bold=True, color=GOLD)

shap_points = [
    ("Top Features (mean |SHAP|)",
     ["1.  age_rating_ratio — dominant driver; age × quality interaction",
      "2.  overall_rating — monotone positive; steeper above rating 80",
      "3.  age — non-linear peak at 24–27, sharp decay after 30",
      "4.  rating_x_potential — interaction captured by boosted trees",
      "5.  goals_per_90 — primary production-rate signal"],
     Inches(0.5), Inches(1.1), Inches(6.0), Inches(4.2)),
    ("Dependence Plots",
     ["age_rating_ratio: smooth monotone increase — validates formula design",
      "overall_rating: non-linear; elbow near rating 78–80",
      "Interaction colour shows age modulates rating SHAP values",
      "Younger high-rated players receive the largest positive SHAP push"],
     Inches(6.8), Inches(1.1), Inches(6.0), Inches(4.2)),
]

for title, bullets, left, top, width, height in shap_points:
    add_bullet_box(sl, title, bullets, left, top, width, height,
                   title_size=16, bullet_size=14)

# Local explanation boxes
add_rect(sl, Inches(0.5), Inches(5.45), Inches(5.8), Inches(1.1), RGBColor(0x0D, 0x2A, 0x15))
add_text(sl, "High-FPVI Player (waterfall)",
         Inches(0.7), Inches(5.5), Inches(5.4), Inches(0.4), size=14, bold=True, color=GREEN)
add_text(sl, "High age_rating_ratio + high rating + strong goals_per_90\n→  prediction +4.8 log-points above base",
         Inches(0.7), Inches(5.9), Inches(5.4), Inches(0.55), size=13, color=LIGHT)

add_rect(sl, Inches(6.7), Inches(5.45), Inches(5.8), Inches(1.1), RGBColor(0x2A, 0x10, 0x10))
add_text(sl, "Low-FPVI Player (waterfall)",
         Inches(6.9), Inches(5.5), Inches(5.4), Inches(0.4), size=14, bold=True, color=YELLOW)
add_text(sl, "Low overall rating + advanced age\n→  prediction pushed sharply below base",
         Inches(6.9), Inches(5.9), Inches(5.4), Inches(0.55), size=13, color=LIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Classification Results
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 8)

add_text(sl, "Transfer Risk Classification", Inches(0.5), Inches(0.3),
         Inches(12), Inches(0.7), size=32, bold=True, color=GOLD)

# Results table
cls_rows = [
    ("Logistic Regression", "45.2%", "0.297", "0.522"),
    ("Random Forest",       "45.0%", "0.316", "0.548"),
    ("XGBoost",             "41.9%", "0.353", "0.542"),
    ("LightGBM",            "41.9%", "0.376", "0.544"),
]
hdrs_c = ["Model", "Accuracy", "Macro F1", "ROC-AUC (OvR)"]
col_xs_c = [Inches(0.5), Inches(4.5), Inches(6.8), Inches(9.2)]
col_ws_c = [Inches(3.9), Inches(2.2), Inches(2.3), Inches(2.4)]

add_text(sl, "Majority-class baseline: 45.2%  (Low risk)",
         Inches(0.5), Inches(1.05), Inches(10.0), Inches(0.4),
         size=14, color=YELLOW, italic=True)

for ci, (h, cx, cw) in enumerate(zip(hdrs_c, col_xs_c, col_ws_c)):
    add_rect(sl, cx, Inches(1.5), cw - Inches(0.05), Inches(0.45),
             RGBColor(0x10, 0x30, 0x50))
    add_text(sl, h, cx + Inches(0.08), Inches(1.55), cw, Inches(0.35),
             size=14, bold=True, color=GOLD)

for ri, row in enumerate(cls_rows):
    top = Inches(1.95) + ri * Inches(0.52)
    bg = ACCENT if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x3E)
    for val, cx, cw in zip(row, col_xs_c, col_ws_c):
        add_rect(sl, cx, top, cw - Inches(0.05), Inches(0.47), bg)
        add_text(sl, val, cx + Inches(0.08), top + Inches(0.05),
                 cw, Inches(0.35), size=13, color=WHITE)

# Interpretation
add_bullet_box(sl, "Interpretation",
    ["All models perform near the majority-class baseline (45%)",
     "ROC-AUC barely above 0.5 — consistent with randomly assigned labels",
     "LightGBM achieves best Macro F1 (0.376) — detects weak non-linear patterns",
     "SHAP still provides domain-consistent insights despite low accuracy"],
    Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.4),
    title_size=16, bullet_size=14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Clustering
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 9)

add_text(sl, "Player Archetype Discovery — K-Means", Inches(0.5), Inches(0.3),
         Inches(12), Inches(0.7), size=32, bold=True, color=GOLD)

# Method box
add_bullet_box(sl, "Method",
    ["K-Means on 10 performance features (standardised)",
     "k selected by silhouette coefficient over k ∈ {2, …, 8}",
     "Optimal k = 2  (silhouette = 0.42)",
     "Visualised with PCA 2D projection + radar chart"],
    Inches(0.5), Inches(1.1), Inches(4.5), Inches(2.5), bullet_size=14)

# Archetype table
add_text(sl, "Recovered Archetypes", Inches(5.3), Inches(1.1),
         Inches(7.5), Inches(0.45), size=16, bold=True, color=GOLD)

arch_rows = [
    ("Elite Players",  "2,229", "28.1", "76.9", "0.70", "0.46", "42.1"),
    ("Goal Scorers",   "  571", "27.5", "76.9", "3.19", "2.51", "75.4"),
]
arch_hdrs = ["Archetype", "n", "Age", "Rating", "G/90", "A/90", "FPVI (M€)"]
arch_xs = [Inches(5.3), Inches(7.35), Inches(8.15), Inches(8.95),
           Inches(9.75), Inches(10.55), Inches(11.4)]
arch_ws = [Inches(1.95), Inches(0.75), Inches(0.75), Inches(0.75),
           Inches(0.75),  Inches(0.80),  Inches(1.1)]

for h, cx, cw in zip(arch_hdrs, arch_xs, arch_ws):
    add_rect(sl, cx, Inches(1.6), cw - Inches(0.04), Inches(0.45),
             RGBColor(0x10, 0x30, 0x50))
    add_text(sl, h, cx + Inches(0.05), Inches(1.65), cw, Inches(0.35),
             size=12, bold=True, color=GOLD)

arch_colors = [RGBColor(0x10, 0x28, 0x45), RGBColor(0x10, 0x2A, 0x18)]
for ri, (row, bg) in enumerate(zip(arch_rows, arch_colors)):
    top = Inches(2.05) + ri * Inches(0.55)
    for val, cx, cw in zip(row, arch_xs, arch_ws):
        add_rect(sl, cx, top, cw - Inches(0.04), Inches(0.5), bg)
        col = GREEN if ri == 1 and val in ["3.19", "2.51", "75.4"] else WHITE
        add_text(sl, val, cx + Inches(0.05), top + Inches(0.06),
                 cw, Inches(0.4), size=12, color=col)

# Key insight
add_rect(sl, Inches(5.3), Inches(3.2), Inches(7.5), Inches(0.7), RGBColor(0x0D, 0x2A, 0x15))
add_text(sl,
    "Goal Scorers: 4.6× higher goals/90 · 5.5× higher assists/90 · +33 M€ FPVI premium",
    Inches(5.5), Inches(3.28), Inches(7.1), Inches(0.5), size=14, bold=True, color=GREEN)

# Insights box
add_bullet_box(sl, "Insights",
    ["K-Means recovers production-based archetypes without position labels",
     "Both clusters share similar age (27–28) and overall rating (~77)",
     "The differentiator is per-90 output — not raw ability",
     "FPVI gap of +33 M€ validates the production weighting in the formula",
     "PCA PC1 (18.3% variance) loads on per-90 stats — separates clusters"],
    Inches(0.5), Inches(3.7), Inches(12.3), Inches(2.8),
    title_size=15, bullet_size=13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusions
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl)
slide_number(sl, 10)

add_rect(sl, 0, 0, Inches(0.25), SLIDE_H, GOLD)

add_text(sl, "Conclusions", Inches(0.55), Inches(0.3), Inches(12), Inches(0.7),
         size=32, bold=True, color=GOLD)

contributions = [
    ("1", "Data Quality Diagnosis",
     "EDA revealed near-zero feature–target correlations in the annotated market values.\n"
     "Responsible ML practice: always verify label quality before modelling.",
     YELLOW),
    ("2", "FPVI Regression  R² = 0.958",
     "LightGBM & XGBoost recover the domain-informed index with R² > 0.95.\n"
     "SHAP confirms the expected age-prime effect and rating non-linearity.",
     GREEN),
    ("3", "Player Archetype Discovery",
     "K-Means (k=2) identifies Elite Players vs Goal Scorers without position labels.\n"
     "Goal Scorers command a +33 M€ FPVI premium driven by per-90 production.",
     RGBColor(0x5D, 0xAD, 0xFF)),
]

for i, (num, title, body, col) in enumerate(contributions):
    top = Inches(1.2) + i * Inches(1.6)
    add_rect(sl, Inches(0.55), top, Inches(0.6), Inches(1.3), col)
    add_text(sl, num, Inches(0.55), top + Inches(0.35), Inches(0.6), Inches(0.6),
             size=28, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.2), top, Inches(11.5), Inches(1.3), ACCENT)
    add_text(sl, title, Inches(1.4), top + Inches(0.05), Inches(11.0), Inches(0.5),
             size=17, bold=True, color=col)
    add_text(sl, body, Inches(1.4), top + Inches(0.55), Inches(11.0), Inches(0.7),
             size=14, color=LIGHT)

# Future work
add_rect(sl, Inches(0.55), Inches(6.1), Inches(11.5), Inches(0.85), RGBColor(0x1A, 0x1A, 0x35))
add_text(sl, "Future Work: ",
         Inches(0.75), Inches(6.18), Inches(1.4), Inches(0.55),
         size=13, bold=True, color=GOLD)
add_text(sl,
    "Real transfer-fee data (Transfermarkt)  ·  Time-series form trajectories  ·  "
    "Multimodal inputs (video stats)  ·  TabNet / attention-based models",
    Inches(2.2), Inches(6.18), Inches(9.6), Inches(0.55), size=13, color=LIGHT)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
prs.save(out_path)
print(f"Saved → {out_path}")
print(f"Slides: {len(prs.slides)}")
